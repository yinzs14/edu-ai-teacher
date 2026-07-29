#!/usr/bin/env python3
"""
Template-based courseware PPT generator.
Reads diagnosis JSON from stdin, optionally loads a teacher-uploaded PPT template,
creates a styled courseware PPT for classroom use.

Usage:
  python generate_courseware_template.py <output.pptx> [<template.pptx>]
  - stdin: JSON with diagnosis data + optional templateStyle
"""
import json
import sys
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Cross-platform font detection ----
def get_font_name():
    """Return best available Chinese font for the current OS."""
    if sys.platform == 'win32':
        return '微软雅黑'
    elif sys.platform == 'darwin':
        return 'PingFang SC'
    else:
        # Linux: prefer Noto Sans CJK
        return 'Noto Sans SC'

FONT = get_font_name()
FONT_TITLE = get_font_name()

# ---- Color palette ----
PRIMARY = RGBColor(0x40, 0x9E, 0xFF)     # Blue
PRIMARY_DARK = RGBColor(0x33, 0x7E, 0xCC)
ACCENT = RGBColor(0x67, 0xC2, 0x3A)       # Green
WARN = RGBColor(0xE6, 0xA2, 0x3C)         # Orange/Yellow
DANGER = RGBColor(0xF5, 0x6C, 0x6C)       # Red
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x30, 0x30, 0x30)
GRAY = RGBColor(0x90, 0x90, 0x90)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=Pt(18),
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name=None):
    """Helper: add a text box with basic formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size if font_size else Pt(18)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name or FONT
    p.alignment = alignment
    return txBox


def add_rich_textbox(slide, left, top, width, height, paragraphs_data,
                     default_size=Pt(16), default_color=BLACK):
    """
    Add a text box with multiple paragraphs.
    paragraphs_data: list of dicts with { text, size, bold, color, alignment }
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, para in enumerate(paragraphs_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = para.get('text', '')
        p.font.size = para.get('size', default_size)
        p.font.bold = para.get('bold', False)
        p.font.color.rgb = para.get('color', default_color)
        p.font.name = para.get('font', FONT)
        p.alignment = para.get('alignment', PP_ALIGN.LEFT)

        if para.get('space_before'):
            p.space_before = para['space_before']
        if para.get('space_after'):
            p.space_after = para['space_after']

    return txBox


def slide_cover(prs, data):
    """Slide 1: Cover — 课件标题"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, PRIMARY)

    student = data.get('studentName', '学生')
    grade = data.get('grade', '')

    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.2),
                f'{student} 个性化学习方案',
                font_size=Pt(44), bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    if grade:
        add_textbox(slide, Inches(0.8), Inches(2.8), Inches(11.7), Inches(0.8),
                    f'适用年级：{grade}',
                    font_size=Pt(22), color=RGBColor(0xBB, 0xDD, 0xFF),
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.8),
                '基于学情诊断 · 精准提升薄弱环节',
                font_size=Pt(18), color=RGBColor(0xBB, 0xDD, 0xFF),
                alignment=PP_ALIGN.CENTER)

    teacher = data.get('teacherName', '')
    if teacher:
        add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.6),
                    f'教师：{teacher}',
                    font_size=Pt(16), color=RGBColor(0x99, 0xCC, 0xFF),
                    alignment=PP_ALIGN.CENTER)


def slide_diagnosis(prs, data):
    """Slide 2: Diagnosis overview — 学情概览"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Title bar
    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
                '学情诊断概览', font_size=Pt(32), bold=True, color=PRIMARY)

    # Divider line
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.6), Inches(1.15), Inches(12), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()

    # Summary
    summary = data.get('summary', '暂无数据')
    add_textbox(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.7),
                f'总体评价：{summary}',
                font_size=Pt(18), color=GRAY)

    # Radar scores - simple text display
    scores = data.get('radarScores', {})
    if scores:
        add_textbox(slide, Inches(0.6), Inches(2.3), Inches(5.5), Inches(0.5),
                    '五维能力评估', font_size=Pt(22), bold=True, color=BLACK)

        y = Inches(3.0)
        dimensions = ['计算', '应用题', '几何', '逻辑', '规律']
        for dim in dimensions:
            score = scores.get(dim, 0)
            bar_color = ACCENT if score >= 60 else (WARN if score >= 40 else DANGER)

            # Label
            add_textbox(slide, Inches(0.8), y, Inches(1.0), Inches(0.4),
                        dim, font_size=Pt(16), color=BLACK)

            # Bar background
            bar_shape = slide.shapes.add_shape(
                1, Inches(2.0), y + Inches(0.05), Inches(6.0), Inches(0.3)
            )
            bar_shape.fill.solid()
            bar_shape.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
            bar_shape.line.fill.background()

            # Bar fill
            fill_width = int(6.0 * score / 100)
            if fill_width > 0:
                fill_shape = slide.shapes.add_shape(
                    1, Inches(2.0), y + Inches(0.05),
                    Inches(min(fill_width, 6.0)), Inches(0.3)
                )
                fill_shape.fill.solid()
                fill_shape.fill.fore_color.rgb = bar_color
                fill_shape.line.fill.background()

            # Score text
            add_textbox(slide, Inches(8.3), y, Inches(0.8), Inches(0.4),
                        str(score), font_size=Pt(16), bold=True, color=bar_color,
                        alignment=PP_ALIGN.RIGHT)

            y += Inches(0.55)

    # Weak points
    weak_points = data.get('weakPoints', [])
    if weak_points:
        y_start = max(y + Inches(0.3), Inches(5.5))
        add_textbox(slide, Inches(0.6), y_start, Inches(5.5), Inches(0.5),
                    '薄弱知识点', font_size=Pt(22), bold=True, color=DANGER)

        y_item = y_start + Inches(0.60)
        for wp in weak_points[:5]:  # max 5 items
            text = f"• {wp.get('name', '未知')}（{wp.get('dimension', '')} {wp.get('score', '')}分）"
            add_textbox(slide, Inches(0.8), y_item, Inches(11.5), Inches(0.4),
                        text, font_size=Pt(15), color=BLACK)
            y_item += Inches(0.45)


def slide_weak_points(prs, data):
    """Slide 3: Detailed weak point analysis with suggestions"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
                '薄弱知识点详解与提升建议', font_size=Pt(32), bold=True, color=PRIMARY)

    shape = slide.shapes.add_shape(
        1, Inches(0.6), Inches(1.15), Inches(12), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()

    weak_points = data.get('weakPoints', [])
    if not weak_points:
        add_textbox(slide, Inches(0.6), Inches(2), Inches(12), Inches(1),
                    '暂无薄弱知识点数据', font_size=Pt(20), color=GRAY)
        return

    y = Inches(1.5)
    for i, wp in enumerate(weak_points[:4]):
        name = wp.get('name', '未知')
        dim = wp.get('dimension', '')
        suggestion = wp.get('suggestion', '')
        score = wp.get('score', 0)

        # Item number badge
        badge = slide.shapes.add_shape(
            1, Inches(0.6), y, Inches(0.5), Inches(0.5)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = DANGER
        badge.line.fill.background()

        add_textbox(slide, Inches(0.6), y, Inches(0.5), Inches(0.5),
                    str(i + 1), font_size=Pt(18), bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER)

        # Title
        add_textbox(slide, Inches(1.3), y, Inches(10.5), Inches(0.4),
                    f'{name}', font_size=Pt(20), bold=True, color=BLACK)

        # Dimension + score
        add_textbox(slide, Inches(1.3), y + Inches(0.4), Inches(10.5), Inches(0.35),
                    f'维度：{dim}   得分：{score}分', font_size=Pt(14), color=GRAY)

        # Suggestion
        if suggestion:
            add_textbox(slide, Inches(1.3), y + Inches(0.85), Inches(10.5), Inches(0.8),
                        f'建议：{suggestion}', font_size=Pt(15), color=BLACK)

        y += Inches(1.5)


def slide_training_plan(prs, data):
    """Slide 4: Training plan — 训练计划"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
                '针对性训练计划', font_size=Pt(32), bold=True, color=PRIMARY)

    shape = slide.shapes.add_shape(
        1, Inches(0.6), Inches(1.15), Inches(12), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()

    # Communication script solutions serve as training plan
    script = data.get('communicationScript', {})
    solutions_text = script.get('solutions', '')
    stage_text = script.get('stageKnowledge', '')

    if stage_text:
        add_textbox(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5),
                    '阶段目标', font_size=Pt(22), bold=True, color=PRIMARY_DARK)
        add_textbox(slide, Inches(0.6), Inches(2.1), Inches(12), Inches(1.0),
                    stage_text, font_size=Pt(16), color=BLACK)

    if solutions_text:
        add_textbox(slide, Inches(0.6), Inches(3.0), Inches(12), Inches(0.5),
                    '提升方案', font_size=Pt(22), bold=True, color=ACCENT)
        add_textbox(slide, Inches(0.6), Inches(3.6), Inches(12), Inches(3.0),
                    solutions_text, font_size=Pt(16), color=BLACK)


def slide_closing(prs, data):
    """Slide 5: Closing — 结尾页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)

    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.0),
                '感谢使用教育AI备课助手',
                font_size=Pt(36), bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    student = data.get('studentName', '学生')
    add_textbox(slide, Inches(0.8), Inches(3.2), Inches(11.7), Inches(0.8),
                f'祝 {student} 学习进步，天天向上！',
                font_size=Pt(24), color=RGBColor(0xBB, 0xDD, 0xFF),
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.6),
                '持续关注学情 · 精准教学辅导',
                font_size=Pt(18), color=RGBColor(0x99, 0xCC, 0xFF),
                alignment=PP_ALIGN.CENTER)


def generate(output_path, data):
    """Main: generate the complete courseware PPT."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_cover(prs, data)
    slide_diagnosis(prs, data)
    slide_weak_points(prs, data)
    slide_training_plan(prs, data)
    slide_closing(prs, data)

    prs.save(output_path)
    return True


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print('Usage: python generate_courseware_template.py <output.pptx>', file=sys.stderr)
        sys.exit(1)

    output_path = sys.argv[1]

    # Read JSON from stdin
    raw = sys.stdin.read()
    data = json.loads(raw)

    try:
        generate(output_path, data)
        print('OK')
    except Exception as e:
        print(f'ERROR: {e}', file=sys.stderr)
        sys.exit(1)
