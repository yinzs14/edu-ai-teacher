#!/usr/bin/env python3
"""
AI备课助手 — 针对性练习PPT生成器
支持 4:3 和 16:9 两种比例
生成讲解用PPT：封面 → 知识点概述 → 逐题讲解 → 总结
"""

import sys
import json
import datetime
import platform
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ==================== 字体兼容 ====================
def _detect_font():
    sys_name = platform.system()
    if sys_name == 'Windows':
        return FONT_NAME
    elif sys_name == 'Darwin':
        return 'PingFang SC'
    else:
        return 'Noto Sans SC'

FONT_NAME = _detect_font()

# ==================== 色彩系统 ====================
TITLE_DARK     = RGBColor(0x1F, 0x29, 0x37)
SUBTITLE_ORANGE = RGBColor(0xF9, 0x73, 0x16)
BODY_GRAY      = RGBColor(0x4B, 0x55, 0x63)
LIGHT_GRAY     = RGBColor(0x6B, 0x72, 0x80)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG        = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT       = RGBColor(0xF9, 0xFA, 0xFB)
ACCENT_BLUE    = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_GREEN   = RGBColor(0x10, 0xB9, 0x81)
ACCENT_RED     = RGBColor(0xEF, 0x44, 0x44)
ACCENT_ORANGE  = RGBColor(0xF9, 0x73, 0x16)
ACCENT_PURPLE  = RGBColor(0x8B, 0x5C, 0xF6)
PHASE_COLORS   = [ACCENT_ORANGE, ACCENT_BLUE, ACCENT_GREEN, ACCENT_PURPLE, ACCENT_RED]

def add_bg_rect(slide, w, h, color=BG_LIGHT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_title_bar(slide, w, title_text, subtitle_text=None, title_size=Pt(28)):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), w - Inches(1.2), Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = title_size
    p.font.bold = True
    p.font.color.rgb = TITLE_DARK
    p.font.name = FONT_NAME

    if subtitle_text:
        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.85), w - Inches(1.2), Inches(0.3))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.color.rgb = SUBTITLE_ORANGE
        p2.font.name = FONT_NAME

    # 底部分隔线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.2), w - Inches(1.2), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    line.line.fill.background()

def make_question_slide(slide, w, h, q_num, total, exercise, color_idx=0):
    """创建单题讲解页"""
    add_bg_rect(slide, w, h)
    color = PHASE_COLORS[color_idx % len(PHASE_COLORS)]

    # 题目编号标签
    label = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.35), Inches(1.8), Inches(0.45)
    )
    label.fill.solid()
    label.fill.fore_color.rgb = color
    label.line.fill.background()
    ltf = label.text_frame
    ltf.word_wrap = True
    lp = ltf.paragraphs[0]
    lp.text = f'第 {q_num}/{total} 题'
    lp.font.size = Pt(16)
    lp.font.bold = True
    lp.font.color.rgb = WHITE
    lp.font.name = FONT_NAME
    lp.alignment = PP_ALIGN.CENTER

    # 题目标题
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), w - Inches(1.0), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = exercise.get('title', f'练习题{q_num}')
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = TITLE_DARK
    p.font.name = FONT_NAME

    # 题目内容
    content = exercise.get('content', '题目内容待补充')
    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), w - Inches(1.6), Inches(3.0))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = content
    p2.font.size = Pt(18)
    p2.font.color.rgb = BODY_GRAY
    p2.font.name = FONT_NAME
    p2.line_spacing = Pt(30)

    # 答案区域（底部）
    answer = exercise.get('answer', '')
    if answer:
        ans_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), h - Inches(2.2), w - Inches(1.0), Inches(1.2)
        )
        ans_card.fill.solid()
        ans_card.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF4)
        ans_card.line.color.rgb = RGBColor(0x86, 0xEF, 0xAC)
        
        ans_tb = slide.shapes.add_textbox(Inches(0.8), h - Inches(2.1), w - Inches(1.6), Inches(1.0))
        atf = ans_tb.text_frame
        atf.word_wrap = True
        ap = atf.paragraphs[0]
        ap.text = f'💡 参考答案：{answer}'
        ap.font.size = Pt(16)
        ap.font.color.rgb = ACCENT_GREEN
        ap.font.name = FONT_NAME

    # 知识点标签
    tags = exercise.get('tags', [])
    if tags:
        y = h - Inches(0.85)
        for i, tag in enumerate(tags[:4]):
            tag_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5 + i * 1.4), y, Inches(1.3), Inches(0.35)
            )
            tag_shape.fill.solid()
            tag_shape.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
            tag_shape.line.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
            ttf = tag_shape.text_frame
            tp = ttf.paragraphs[0]
            tp.text = tag
            tp.font.size = Pt(10)
            tp.font.color.rgb = ACCENT_BLUE
            tp.font.name = FONT_NAME
            tp.alignment = PP_ALIGN.CENTER

    # 页脚
    ft = slide.shapes.add_textbox(Inches(0.5), h - Inches(0.45), w - Inches(1.0), Inches(0.35))
    fp = ft.text_frame.paragraphs[0]
    fp.text = f'AI备课助手 · 针对性练习'
    fp.font.size = Pt(8)
    fp.font.color.rgb = LIGHT_GRAY
    fp.font.name = FONT_NAME
    fp.alignment = PP_ALIGN.RIGHT


def generate(data, output_path):
    """生成针对性练习PPT"""
    aspect = data.get('aspect', '16:9')  # '4:3' or '16:9'
    
    if aspect == '4:3':
        w = Inches(10)
        h = Inches(7.5)
    else:
        w = Inches(13.333)
        h = Inches(7.5)

    prs = Presentation()
    prs.slide_width = w
    prs.slide_height = h

    student = data.get('studentName', '同学')
    teacher = data.get('teacherName', '老师')
    weak_point = data.get('weakPoint', {})
    wp_name = weak_point.get('name', '综合练习')
    wp_suggestion = weak_point.get('suggestion', '')
    exercises = data.get('exercises', [])

    # ===== 封面 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, w, h, TITLE_DARK)

    # 标题
    tb = slide.shapes.add_textbox(Inches(0.8), h * 0.25, w - Inches(1.6), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'{wp_name} 针对性练习'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME

    # 副标题
    tb2 = slide.shapes.add_textbox(Inches(0.8), h * 0.48, w - Inches(1.6), Inches(0.8))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f'教师：{teacher}  |  学生：{student}'
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
    p2.font.name = FONT_NAME
    p3 = tf2.add_paragraph()
    p3.text = f'共 {len(exercises)} 道针对性练习题  |  {datetime.date.today().strftime("%Y.%m.%d")}'
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)
    p3.font.name = FONT_NAME

    if wp_suggestion:
        tb3 = slide.shapes.add_textbox(Inches(1.5), h * 0.65, w - Inches(3.0), Inches(1.2))
        tf3 = tb3.text_frame
        tf3.word_wrap = True
        p4 = tf3.paragraphs[0]
        p4.text = f'📋 训练目标：{wp_suggestion[:150]}'
        p4.font.size = Pt(14)
        p4.font.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
        p4.font.name = FONT_NAME

    # ===== 知识点概述页 =====
    if len(exercises) > 1:
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg_rect(slide2, w, h)
        add_title_bar(slide2, w, '练习总览', f'共 {len(exercises)} 道题 · 由易到难递进')

        # 题目列表
        card_w = w - Inches(2.0)
        card_h = Inches(0.65)
        x0 = Inches(1.0)
        start_y = Inches(1.6)
        gap = Inches(0.08)

        for i, ex in enumerate(exercises[:8]):
            color = PHASE_COLORS[i % len(PHASE_COLORS)]
            top = start_y + (card_h + gap) * i
            
            # 编号圆点
            dot = slide2.shapes.add_shape(
                MSO_SHAPE.OVAL, x0, top + Inches(0.15), Inches(0.3), Inches(0.3)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()
            dtf = dot.text_frame
            dp = dtf.paragraphs[0]
            dp.text = str(i + 1)
            dp.font.size = Pt(11)
            dp.font.color.rgb = WHITE
            dp.font.name = FONT_NAME
            dp.alignment = PP_ALIGN.CENTER

            # 标题文字
            ttb = slide2.shapes.add_textbox(x0 + Inches(0.5), top + Inches(0.08), card_w - Inches(0.5), Inches(0.5))
            tf = ttb.text_frame
            tp = tf.paragraphs[0]
            tp.text = ex.get('title', f'练习题{i+1}')
            tp.font.size = Pt(15)
            tp.font.color.rgb = TITLE_DARK
            tp.font.name = FONT_NAME

        # 页脚
        ft = slide2.shapes.add_textbox(Inches(0.5), h - Inches(0.45), w - Inches(1.0), Inches(0.35))
        fp = ft.paragraphs[0]
        fp.text = '建议逐题讲解，每题控制在 5-8 分钟'
        fp.font.size = Pt(11)
        fp.font.color.rgb = LIGHT_GRAY
        fp.font.name = FONT_NAME
        fp.alignment = PP_ALIGN.CENTER

    # ===== 逐题讲解 =====
    for i, ex in enumerate(exercises):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        make_question_slide(s, w, h, i + 1, len(exercises), ex, i)

    # ===== 结尾页 =====
    slide_end = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide_end, w, h, TITLE_DARK)
    
    etb = slide_end.shapes.add_textbox(Inches(0.8), h * 0.3, w - Inches(1.6), Inches(1.0))
    etf = etb.text_frame
    ep = etf.paragraphs[0]
    ep.text = '练习完成！'
    ep.font.size = Pt(40)
    ep.font.bold = True
    ep.font.color.rgb = WHITE
    ep.font.name = FONT_NAME
    ep.alignment = PP_ALIGN.CENTER

    etb2 = slide_end.shapes.add_textbox(Inches(0.8), h * 0.45, w - Inches(1.6), Inches(0.6))
    etf2 = etb2.text_frame
    ep2 = etf2.paragraphs[0]
    ep2.text = f'建议整理错题本，重点标记 {wp_name} 相关内容'
    ep2.font.size = Pt(18)
    ep2.font.color.rgb = LIGHT_GRAY
    ep2.font.name = FONT_NAME
    ep2.alignment = PP_ALIGN.CENTER

    prs.save(output_path)
    return output_path


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print('Usage: python generate_practice.py <output_path>', file=sys.stderr)
        sys.exit(1)

    raw = sys.stdin.read()
    if not raw or not raw.strip():
        print('Error: empty input data', file=sys.stderr)
        sys.exit(1)

    try:
        data = json.loads(raw)
    except json.JSONDecodeError as e:
        print(f'Error: invalid JSON input: {e}', file=sys.stderr)
        sys.exit(1)

    if not data.get('exercises'):
        print('Warning: no exercises data, generating with empty content', file=sys.stderr)

    output = sys.argv[1]
    try:
        result = generate(data, output)
        print(result)
    except Exception as e:
        print(f'Error generating practice PPT: {e}', file=sys.stderr)
        sys.exit(1)
