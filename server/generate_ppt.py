#!/usr/bin/env python3
"""
AI备课助手 — PPT生成器
基于「数学冲刺学习计划」模板风格，用python-pptx动态生成学习方案PPT。
输入：JSON格式诊断数据（通过stdin）
输出：pptx文件路径
"""

import sys
import json
import os
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ==================== 色彩系统（与模板保持一致） ====================
TITLE_DARK   = RGBColor(0x1F, 0x29, 0x37)
SUBTITLE_ORANGE = RGBColor(0xF9, 0x73, 0x16)
BODY_GRAY    = RGBColor(0x4B, 0x55, 0x63)
LIGHT_GRAY   = RGBColor(0x6B, 0x72, 0x80)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG      = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT     = RGBColor(0xF9, 0xFA, 0xFB)
BAR_ORANGE   = RGBColor(0xF9, 0x73, 0x16)
DEEP_ORANGE  = RGBColor(0xC2, 0x41, 0x0C)
ACCENT_BLUE  = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_RED   = RGBColor(0xEF, 0x44, 0x44)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
PHASE_COLORS = [BAR_ORANGE, ACCENT_BLUE, ACCENT_GREEN]

# ==================== 工具函数 ====================

def add_bg_rect(slide, color=BG_LIGHT):
    """给slide添加纯色背景矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_title_bar(slide, title_text, subtitle_text=None, title_size=Pt(32)):
    """添加统一的标题区"""
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = title_size
    p.font.bold = True
    p.font.color.rgb = TITLE_DARK
    p.font.name = '微软雅黑'

    if subtitle_text:
        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12.133), Inches(0.35))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(18)
        p2.font.color.rgb = SUBTITLE_ORANGE
        p2.font.name = '微软雅黑'

def make_card(slide, left, top, width, height, bar_color, title, lines, title_size=Pt(14), body_size=Pt(12)):
    """创建带色条顶部的卡片"""
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    card.line.width = Pt(0.5)

    # 顶部色条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.07)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()

    # 标题
    tb = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.35)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = title_size
    p.font.bold = True
    p.font.color.rgb = TITLE_DARK
    p.font.name = '微软雅黑'

    # 内容
    tb2 = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.55), width - Inches(0.3), height - Inches(0.7)
    )
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = line
        p.font.size = body_size
        p.font.color.rgb = BODY_GRAY
        p.font.name = '微软雅黑'
        p.space_after = Pt(4)

def add_speaker_notes(slide, text):
    """添加演讲者备注"""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text

# ==================== 各页面构建 ====================

def slide_cover(prs, data):
    """封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg_rect(slide, TITLE_DARK)

    # 标题
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.733), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    subject = data.get('subject', '数学')
    p.text = f'{subject}冲刺学习计划'
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = '微软雅黑'

    # 教师姓名
    teacher = data.get('teacherName', '老师')
    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(6), Inches(1.0))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f'授课教师：{teacher}'
    p2.font.size = Pt(18)
    p2.font.color.rgb = WHITE
    p2.font.name = '微软雅黑'
    p3 = tf2.add_paragraph()
    p3.text = datetime.date.today().strftime('%Y.%m.%d')
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
    p3.font.name = '微软雅黑'

    add_speaker_notes(slide, f'开场介绍：这是为{data.get("studentName", "同学")}定制的{subject}学习计划。先讲目标和整体安排，建立家长信心。')

def slide_diagnosis(prs, data):
    """学情诊断结果页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide)
    student = data.get('studentName', '同学')
    add_title_bar(slide, '学情诊断结果', f'根据{student}近期作业/试卷分析 · 明确薄弱方向')

    scores = data.get('radarScores', {})
    dims = list(scores.keys())
    start_y = Inches(1.6)
    bar_w = Inches(7.5)
    bar_h = Inches(0.28)
    gap = Inches(0.5)

    for i, dim in enumerate(dims[:5]):
        score = scores[dim]
        y = start_y + gap * i

        # 维度标签
        tb = slide.shapes.add_textbox(Inches(0.8), y, Inches(1.6), bar_h)
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = dim
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_DARK
        p.font.name = '微软雅黑'

        # 底色条
        bg_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), y + Inches(0.02), bar_w, bar_h - Inches(0.04)
        )
        bg_bar.fill.solid()
        bg_bar.fill.fore_color.rgb = RGBColor(0xF3, 0xF4, 0xF6)
        bg_bar.line.fill.background()

        # 分数条
        if score >= 70:
            fill_c = ACCENT_GREEN
        elif score >= 50:
            fill_c = BAR_ORANGE
        else:
            fill_c = ACCENT_RED
        score_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), y + Inches(0.02),
            bar_w * (score / 100), bar_h - Inches(0.04)
        )
        score_bar.fill.solid()
        score_bar.fill.fore_color.rgb = fill_c
        score_bar.line.fill.background()

        # 分数
        tb2 = slide.shapes.add_textbox(Inches(10.2), y, Inches(0.8), bar_h)
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f'{score}分'
        p2.font.size = Pt(13)
        p2.font.color.rgb = LIGHT_GRAY
        p2.font.name = '微软雅黑'
        p2.alignment = PP_ALIGN.RIGHT

    # 总结
    summary = data.get('summary', '')
    if summary:
        tb3 = slide.shapes.add_textbox(Inches(0.8), start_y + gap * len(dims) + Inches(0.3), Inches(12), Inches(0.6))
        tf3 = tb3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = f'📊 {summary}'
        p3.font.size = Pt(14)
        p3.font.color.rgb = LIGHT_GRAY
        p3.font.name = '微软雅黑'

    add_speaker_notes(slide, '先问孩子觉得自己哪里最弱，再展示数据。每个维度用简单语言解释分数含义。')

def slide_weak_points(prs, data):
    """薄弱知识点分析页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide)
    add_title_bar(slide, '薄弱知识点分析', 'Core Weakness Analysis — 聚焦问题，精准提分')

    points = data.get('weakPoints', [])
    if not points:
        points = [{'name': '暂未识别', 'dimension': '请先完成学情诊断', 'score': 0, 'suggestion': '先上传作业进行诊断'}]

    max_show = 6
    display = points[:max_show]
    cols = min(3, len(display))
    card_w = Inches(3.7)
    card_h = Inches(2.4)
    x0 = Inches(0.6)
    gap_x = Inches(0.25)
    start_y = Inches(1.6)
    gap_y = Inches(0.2)

    for i, wp in enumerate(display):
        col = i % cols
        row = i // cols
        left = x0 + (card_w + gap_x) * col
        top = start_y + (card_h + gap_y) * row

        # 卡片
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)

        # 顶部分数条
        score = wp.get('score', 0)
        zhanbi = max(5, score)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w * (zhanbi / 100), Inches(0.06)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT_RED if score < 50 else BAR_ORANGE
        bar.line.fill.background()

        # 知识点名
        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), card_w - Inches(0.4), Inches(0.35))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = wp.get('name', '未知')
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = TITLE_DARK
        p.font.name = '微软雅黑'

        # 维度标签
        tb2 = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.55), card_w - Inches(0.4), Inches(0.25))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f'[{wp.get("dimension", "")}] 掌握度：{score}分'
        p2.font.size = Pt(11)
        p2.font.color.rgb = LIGHT_GRAY
        p2.font.name = '微软雅黑'

        # 建议
        tb3 = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.0), card_w - Inches(0.4), Inches(1.1))
        tf3 = tb3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = wp.get('suggestion', '')
        p3.font.size = Pt(12)
        p3.font.color.rgb = BODY_GRAY
        p3.font.name = '微软雅黑'

    add_speaker_notes(slide, '逐项和孩子确认每个知识点的掌握程度。用问句引导：这个类型的题目做起来顺手吗？')

def slide_plan_overview(prs, data):
    """学习方案总览页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide)
    add_title_bar(slide, '备课方案：三阶段冲刺计划', '3-Stage Sprint Plan')

    # 顶部3个总览卡片
    overviews = [
        ('循序渐进', '从开卷到闭卷，从模块到综合'),
        ('每周频次', '每周3-5节课，每次1-2小时'),
        ('核心方法', '真题导向，反复练习，整理归纳'),
    ]
    card_w = Inches(3.7)
    x0 = Inches(0.6)
    start_y = Inches(1.6)
    for i, (title, desc) in enumerate(overviews):
        left = x0 + (card_w + Inches(0.25)) * i
        top = start_y
        make_card(slide, left, top, card_w, Inches(0.9), PHASE_COLORS[i % 3], title, [desc], Pt(14), Pt(12))

    # 三阶段卡片
    phases = build_phases(data)
    phase_y = Inches(2.8)
    for i, ph in enumerate(phases[:3]):
        left = x0 + (card_w + Inches(0.25)) * i
        make_card(slide, left, phase_y, card_w, Inches(2.5), PHASE_COLORS[i],
                  f'0{i+1} {ph["title"]}', ph['lines'], Pt(16), Pt(12))

    add_speaker_notes(slide, '用三阶段法帮孩子建立信心：第一阶段不要求速度，第二阶段不要求全对，第三阶段才是真正的冲刺。')

def slide_phase_detail(prs, data, phase_idx):
    """阶段详情页（3个阶段分别调用）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide)

    phases = build_phases(data)
    if phase_idx >= len(phases):
        return
    ph = phases[phase_idx]

    en_titles = ['Stage 1: Open-Book Familiarization', 'Stage 2: Closed-Book Consolidation', 'Stage 3: Comprehensive Review']
    en_title = en_titles[phase_idx] if phase_idx < len(en_titles) else ''

    add_title_bar(slide, f'第{phase_idx+1}阶段：{ph["title"]}', en_title)

    # 三栏布局
    col_w = Inches(3.7)
    x0 = Inches(0.6)
    y0 = Inches(1.6)
    card_h = Inches(4.5)

    # 左栏：阶段目标
    make_card(slide, x0, y0, col_w, card_h, PHASE_COLORS[phase_idx], '阶段目标',
              ph['goals'], Pt(18), Pt(13))

    # 中栏：核心方法
    make_card(slide, x0 + col_w + Inches(0.25), y0, col_w, card_h, PHASE_COLORS[phase_idx], '核心方法',
              ph['methods'], Pt(18), Pt(13))

    # 右栏：时间规划
    right = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x0 + (col_w + Inches(0.25)) * 2, y0, col_w, card_h
    )
    right.fill.solid()
    right.fill.fore_color.rgb = WHITE
    right.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)

    tb = slide.shapes.add_textbox(x0 + (col_w + Inches(0.25)) * 2 + Inches(0.3), y0 + Inches(0.8), col_w - Inches(0.6), Inches(1.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = ph['time']
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = SUBTITLE_ORANGE
    p.font.name = '微软雅黑'
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = ph['classes']
    p2.font.size = Pt(16)
    p2.font.color.rgb = BODY_GRAY
    p2.font.name = '微软雅黑'
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = '\n建议每天安排固定时间\n循序渐进完成模块拆解'
    p3.font.size = Pt(13)
    p3.font.color.rgb = LIGHT_GRAY
    p3.font.name = '微软雅黑'
    p3.alignment = PP_ALIGN.CENTER

    add_speaker_notes(slide, f'{ph["title"]}：详细介绍目标和具体方法，让孩子和家长都清楚这个阶段要做什么、怎么做。')

def slide_schedule(prs, data):
    """学习安排页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide)
    add_title_bar(slide, '第一阶段：学习安排', 'Weekly Schedule (Stage 1)')

    points = data.get('weakPoints', [])
    if not points:
        points = [{'name': '基础概念', 'dimension': '计算'}]

    # 为每个薄弱点生成一节课
    cards = []
    for i, wp in enumerate(points[:8]):
        date_str = f'第{i+1}课'
        cards.append({
            'title': wp.get('name', f'知识点{i+1}'),
            'desc': wp.get('suggestion', '系统学习，夯实基础')[:40],
            'color': PHASE_COLORS[i % 3]
        })

    cols = 4
    card_w = Inches(2.85)
    card_h = Inches(1.5)
    x0 = Inches(0.4)
    gap_x = Inches(0.15)
    start_y = Inches(1.6)
    gap_y = Inches(0.15)

    for i, card_data in enumerate(cards):
        col = i % cols
        row = i // cols
        left = x0 + (card_w + gap_x) * col
        top = start_y + (card_h + gap_y) * row

        c = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h
        )
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)

        # 左侧色条
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, Inches(0.05), card_h
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = card_data['color']
        bar.line.fill.background()

        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), card_w - Inches(0.3), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = card_data['title']
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = SUBTITLE_ORANGE
        p.font.name = '微软雅黑'

        tb2 = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.65), card_w - Inches(0.3), Inches(0.7))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = card_data['desc']
        p2.font.size = Pt(12)
        p2.font.color.rgb = BODY_GRAY
        p2.font.name = '微软雅黑'

    # 底部提示
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(5.95), Inches(12), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = '进度根据个人掌握情况灵活调整，稳扎稳打最重要'
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = '微软雅黑'
    p.alignment = PP_ALIGN.CENTER

    add_speaker_notes(slide, '课表可根据实际情况灵活调整。每节课聚焦一个知识点，不贪多，稳扎稳打。')

def slide_communication(prs, data):
    """家长沟通话术页"""
    script = data.get('communicationScript', {}) or {}
    has_content = any(script.get(k) for k in ['stageKnowledge', 'mastered', 'weaknesses', 'solutions', 'talkingTips'])
    if not has_content:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide)
    student = data.get('studentName', '同学')
    add_title_bar(slide, '家长沟通话术', f'可直接口述给 {student} 家长的专业沟通脚本')

    sections = [
        ('本阶段知识', 'stageKnowledge', ACCENT_PURPLE, '📖'),
        ('已掌握部分', 'mastered', ACCENT_GREEN, '✅'),
        ('有待提升', 'weaknesses', ACCENT_RED, '⚠️'),
        ('解决建议', 'solutions', ACCENT_BLUE, '💡'),
        ('沟通要点', 'talkingTips', LIGHT_GRAY, '💬'),
    ]

    active_sections = [(label, key, color, icon) for (label, key, color, icon) in sections if script.get(key)]

    if not active_sections:
        return

    card_h = Inches(1.4)
    gap = Inches(0.12)
    max_h = Inches(5.2)
    # 动态计算卡片高度
    per_card = min(card_h, (max_h - gap * (len(active_sections) - 1)) / len(active_sections))
    y0 = Inches(1.55)

    for i, (label, key, color, icon) in enumerate(active_sections):
        top = y0 + (per_card + gap) * i
        content = script.get(key, '')
        preview = content[:200] + ('…' if len(content) > 200 else '')

        # 卡片
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), top, Inches(12.133), per_card
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)

        # 左侧色条
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(0.06), per_card
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # 标签
        tb = slide.shapes.add_textbox(Inches(0.85), top + Inches(0.08), Inches(2.5), Inches(0.3))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = f'{icon} {label}'
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = '微软雅黑'

        # 内容
        tb2 = slide.shapes.add_textbox(Inches(0.85), top + Inches(0.4), Inches(11.3), per_card - Inches(0.5))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = preview
        p2.font.size = Pt(11)
        p2.font.color.rgb = BODY_GRAY
        p2.font.name = '微软雅黑'

    add_speaker_notes(slide, '此页展示了可以直接念给家长听的沟通话术，按顺序读即可。')

def slide_closing(prs, data):
    """结尾页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, TITLE_DARK)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.733), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = '你的坚持，终将美好'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = '微软雅黑'
    p.alignment = PP_ALIGN.CENTER

    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.733), Inches(0.5))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = 'Your Persistence Will Pay Off'
    p2.font.size = Pt(18)
    p2.font.color.rgb = LIGHT_GRAY
    p2.font.name = '微软雅黑'
    p2.alignment = PP_ALIGN.CENTER

    # 名言
    tb3 = slide.shapes.add_textbox(Inches(1.5), Inches(2.6), Inches(10.333), Inches(1.0))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = '"我们的目标不是成为数学天才，而是成为一个高效的得分手。"'
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.font.color.rgb = DEEP_ORANGE
    p3.font.name = '微软雅黑'
    p3.alignment = PP_ALIGN.CENTER

    # 鼓励语
    items = [
        ('相信自己', '你比想象中更强大'),
        ('紧跟计划', '每一步都算数'),
        ('保持专注', '把精力放在能得分的地方'),
        ('永不言弃', '坚持到最后一刻'),
    ]
    card_w = Inches(2.7)
    x0 = Inches(0.8)
    y0 = Inches(4.2)
    for i, (title, desc) in enumerate(items):
        left = x0 + (card_w + Inches(0.2)) * i
        make_card(slide, left, y0, card_w, Inches(1.2), PHASE_COLORS[i % 3], title, [desc], Pt(14), Pt(11))

# ==================== 数据处理 ====================

def build_phases(data):
    """根据诊断数据构建三阶段"""
    points = data.get('weakPoints', [])
    p1_names = [wp.get('name', '') for wp in points[:3]]
    p1_str = '、'.join(p1_names) if p1_names else '基础知识'

    # 用AI生成的描述，如果没有则用默认值
    script = data.get('communicationScript', {}) or {}
    solutions = script.get('solutions', '')
    weaknesses = script.get('weaknesses', '')

    return [
        {
            'title': '开卷熟悉（识别与应用）',
            'lines': [
                f'重点攻克：{p1_str}',
                '认识题型，知道用什么知识',
                '允许查公式、看笔记，独立思考完成',
            ],
            'goals': [
                '解决"识别"问题',
                '快速判断题目所属知识模块',
                '建立"题目—模块"条件反射',
                '准确匹配公式与定理',
            ],
            'methods': [
                '开卷模式：允许查阅资料',
                '独立思考：先尝试后验证',
                '节奏灵活：基础好2个/天',
                '薄弱模块放慢到1个/天',
            ],
            'time': '约 1-2 周',
            'classes': '（共计3-5节课）',
        },
        {
            'title': '闭卷巩固（背诵与记忆）',
            'lines': [
                f'检验前阶段学习成果',
                '脱离资料，独立完成',
                '限时训练，整理错题',
            ],
            'goals': [
                '解决"背诵"问题',
                '核心公式能脱离资料默写',
                '检验学习成果',
                '独立完成基础题目',
            ],
            'methods': [
                '闭卷模式：限时完成题目',
                '整理公式本：汇总必考公式',
                '错题复盘：记录分析原因',
                '避免重复犯错',
            ],
            'time': '约 1 周',
            'classes': '（共计2-3节课）',
        },
        {
            'title': '总复习（综合与提升）',
            'lines': [
                '真题套练，模拟考试',
                '查漏补缺，弥补漏洞',
                '保持良好心态迎考',
            ],
            'goals': [
                '适应考试节奏',
                '提高做题速度和准确率',
                '查漏补缺',
                '从容自信迎接挑战',
            ],
            'methods': [
                '真题套练：近3年真题模拟',
                '模块复习：针对性强化',
                '错题复盘：回顾前两阶段',
                '避免重复踩坑',
            ],
            'time': '课程剩余时间',
            'classes': '（直至考试）',
        },
    ]

# ==================== 主入口 ====================

def generate(data, output_path):
    """生成PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    prs.core_properties.author = data.get('teacherName', 'AI备课助手')
    prs.core_properties.title = f'{data.get("studentName", "学生")} {data.get("subject", "数学")} 学习方案'

    # 按模板顺序生成所有页面
    slide_cover(prs, data)
    slide_diagnosis(prs, data)
    slide_weak_points(prs, data)
    slide_plan_overview(prs, data)
    slide_phase_detail(prs, data, 0)
    slide_phase_detail(prs, data, 1)
    slide_phase_detail(prs, data, 2)
    slide_schedule(prs, data)
    slide_communication(prs, data)
    slide_closing(prs, data)

    prs.save(output_path)
    return output_path

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print('Usage: python generate_ppt.py <output_path>', file=sys.stderr)
        print('Input: JSON data from stdin', file=sys.stderr)
        sys.exit(1)

    raw = sys.stdin.read()
    data = json.loads(raw)
    output = sys.argv[1]
    result = generate(data, output)
    print(result)
